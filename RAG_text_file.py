from llama_index.core import VectorStoreIndex, SimpleDirectoryReader
from llama_index.core.agent.workflow import FunctionAgent
import chromadb
from llama_index.vector_stores.chroma import ChromaVectorStore
from llama_index.core import StorageContext

import openai
from llama_index.llms.openai import OpenAI
import asyncio
from dotenv import load_dotenv
import os
from llama_index.core.response.pprint_utils import pprint_response
import PyPDF2
from fastapi import FastAPI, HTTPException, Request
from fastapi.middleware.cors import CORSMiddleware
import uvicorn
import logging
from typing import List, Dict, Any


######################################################################################
from llama_index.core import Settings

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Configuration
CONFIG = {
    "DIRECTORY_PATH": "D:\\AppBulls\\LLM Application Frameworks\\llama_index\\llamaIndex_modules\\Directory_Summarizer\\documents",
    "file_extensions": ['.docx', '.pptx', '.txt', '.pdf'],
    "chunk_size": 512,
    "chunk_overlap": 64,
    "similarity_top_k": 5,
    "model_name": "gpt-4o",
    "temperature": 0.3
}

# Local settings
from llama_index.core.node_parser import SentenceSplitter

###################################################################################

#I want to use a different vector store


load_dotenv()
openai.api_key = os.getenv("OPENAI_API_KEY")

# Initialize components
def initialize_components():
    """Initialize all required components with error handling"""
# Check if directory exists
    try:
        if not os.path.exists(CONFIG["DIRECTORY_PATH"]):
            raise FileNotFoundError(f"The directory {CONFIG['DIRECTORY_PATH']} does not exist")

    #RAG implementation
        print("Files in directory:", os.listdir(CONFIG['DIRECTORY_PATH']))
    
        chroma_client = chromadb.PersistentClient()
        chroma_collection = chroma_client.get_or_create_collection("quickstart")
        vector_store = ChromaVectorStore(chroma_collection=chroma_collection)
        storage_context = StorageContext.from_defaults(vector_store=vector_store)   

        Settings.chunk_size = CONFIG['chunk_size']
    # Load documents with error handling
        logger.info("Loading documents...")
        documents = SimpleDirectoryReader(CONFIG['DIRECTORY_PATH'], required_exts=CONFIG['file_extensions']).load_data()
        index = VectorStoreIndex.from_documents(documents, storage_context=storage_context, transformations=[SentenceSplitter(chunk_overlap = CONFIG['chunk_overlap'], chunk_size=CONFIG['chunk_size'])])
        query_engine = index.as_query_engine(Streaming = True, similarity_top_k = 5, chat_mode="best", verbose = False)

        return query_engine

    except Exception as e:
        logger.error(f"Initialization failed: {str(e)}")
        raise

from docx import Document
# File reading functions
def read_txt(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        return f.read()
    
def read_pdf(file_path):
    content = []
    with open(file_path, 'rb') as f:
        reader = PyPDF2.PdfReader(f)
        for page in reader.pages:
            text = page.extract_text()
            if text:
                content.append(text)
    return "\n".join(content)

# Read .docx files
from docx import Document
def read_docx(file_path):
    try:
        doc = Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
    except Exception as e:
        return f"Error reading docs document{file_path}: {e}"
    
# Read .pptx files
from pptx import Presentation
def read_pptx(file_path):
    try:
        prs = Presentation(file_path)
        content = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    content.append(shape.text)
        return "\n".join(content)
    
    except Exception as e:
        return f"Error reading pptx document{file_path}: {e}"
    
    # Tool 3: Read the content of a given file
def read_file(file_name: str, base_path: str = CONFIG['DIRECTORY_PATH']) -> str:
    file_name = file_name.strip()
    # Construct the expected full path
    full_path = os.path.join(base_path, file_name)
    
    # If the file is not found by exact name, search for close matches
    if not os.path.isfile(full_path):
        # List all files in the directory
        available_files = os.listdir(base_path)
        # Look for files that contain the given file_name (case-insensitive)
        matching_files = [f for f in available_files if file_name.lower() in f.lower()]
        
        if matching_files:
            # Optionally, print/log the found files
            print(f"File '{file_name}' not found exactly. Using closest match: {matching_files[0]}")
            full_path = os.path.join(base_path, matching_files[0])
        else:
            return f"{file_name} is not a valid file."
    
    try:
        ext = os.path.splitext(full_path)[1].lower()
        if ext == ".docx":
            return read_docx(full_path)
        elif ext == ".pptx":
            return read_pptx(full_path)
        elif ext == ".txt":
            return read_txt(full_path)
        elif ext == ".pdf":
            return read_pdf(full_path)
        else:
            return f"Unsupported file format: {ext}"
    except Exception as e:
        return f"Error reading {file_name}: {e}"
    
# Tool: Search documents using the query engine
try:
    query_engine = initialize_components()
    async def search_document(query):
        try:
            response = await query_engine.aquery(query)
            return str(response)
        except Exception as e:
            logger.error(f"Search error: {str(e)}")
            return f"Search error: {str(e)}"


# Tool to list all files and folders in a directory
    def list_directory(path: str = CONFIG['DIRECTORY_PATH']) -> List[str]:
        entries = os.listdir(path)
        return entries


    agent = FunctionAgent(
    tools=[search_document, list_directory, read_file],
    llm = OpenAI(model = CONFIG['model_name']),
    temperature = CONFIG['temperature'],
    # top_p = 1,
    system_prompt = 
    """
        You are an intelligent job finder system that helps match candidates with startups.
        When processing queries:
        - Search indexed documents (search_document)
        - List available files (list_directory)
        - Read file contents when needed (read_file)
        
        Guidelines:
        1. Analyze resumes and job descriptions carefully
        2. Match skills and experiences accurately
        3. Provide clear, concise recommendations
        4. Ignore non-text elements like images
        5. Be professional and helpful
"""
) 
except Exception as e:
    logger.critical(f"Failed to initialize application: {str(e)}")
    raise

#############Initialize fast api#########################
# Initialize FastAPI app (unchanged)
app = FastAPI()

# Enable CORS (unchanged)
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:8000"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)
@app.post("/job_search")
async def query_database(request: Request):
    try:
        data = await request.json()
        query = data.get("query", "").strip()

        if not query:
            raise HTTPException(400, "Query cannot be empty")
    
        logger.info(f"Processing query: {query}")

        if query.lower() in ["bye", "exit", "close"]:
            return {"response": "See you later!"}

        try:
            response = await agent.run(query)
            return {"response": str(response)}
        except Exception as e:
            logger.error(f"Query processing error: {str(e)}")
            raise HTTPException(status_code=500, detail="Error processing your query")
    
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Unexpected error: {str(e)}")
        raise HTTPException(status_code=500, detail="Internal server error")

if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=8000)